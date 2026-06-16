from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x1B, 0x5E, 0x20)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF1, 0xF8, 0xE9)
GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0x2E, 0x7D, 0x32)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=GRAY, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_bullet_slide(slide, title, items, sub_title=""):
    add_bg(slide, WHITE)
    # Top bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)

    if sub_title:
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12), Inches(0.5))
        tf2 = txBox.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = sub_title
        p2.font.size = Pt(16)
        p2.font.italic = True
        p2.font.color.rgb = ACCENT
        p2.font.name = "Calibri"

    y_start = Inches(2.0) if sub_title else Inches(1.5)
    txBox = slide.shapes.add_textbox(Inches(0.6), y_start, Inches(12), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(6)
        p.level = 0

def add_two_col_slide(slide, title, left_items, right_items, left_title="", right_title=""):
    add_bg(slide, WHITE)
    # Top bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)

    # Left column header
    if left_title:
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.5))
        tf2 = txBox.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = left_title
        p2.font.size = Pt(20)
        p2.font.bold = True
        p2.font.color.rgb = ACCENT
        p2.font.name = "Calibri"

    # Left column
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(5.8), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(17)
        p.font.color.rgb = GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(5)

    # Right column header
    if right_title:
        txBox = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(5.8), Inches(0.5))
        tf2 = txBox.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = right_title
        p2.font.size = Pt(20)
        p2.font.bold = True
        p2.font.color.rgb = ACCENT
        p2.font.name = "Calibri"

    # Right column
    txBox = slide.shapes.add_textbox(Inches(6.8), Inches(2.0), Inches(5.8), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(17)
        p.font.color.rgb = GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(5)


# ═══════════════════ TITLE SLIDE ═══════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
    "TagX", font_size=60, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.0), Inches(11), Inches(1),
    "Business Model Canvas", font_size=36, bold=False, color=RGBColor(0xA5, 0xD6, 0xA7), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(4.2), Inches(11), Inches(1),
    "India's Smart Tracking Tag — Hardware + Subscription Platform", font_size=20, bold=False, color=RGBColor(0xC8, 0xE6, 0xC9), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
    "Navi Mumbai, India  |  2026–2030", font_size=18, bold=False, color=RGBColor(0x81, 0xC7, 0x84), alignment=PP_ALIGN.CENTER)

# ═══════════════════ 1. CUSTOMER SEGMENTS ═══════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_two_col_slide(slide, "1. Customer Segments",
    [
        "● B2C — Indian Consumers (Mass Market)",
        "   • Individuals tracking keys, wallets, bags",
        "   • Families tracking pets, children, luggage",
        "   • Vehicle owners (two-wheeler, car)",
        "   • Age 18–45, urban & semi-urban India",
        "",
        "● B2B — Indian Enterprises (Niche)",
        "   • Small fleet operators (10–200 vehicles)",
        "   • Logistics & delivery companies",
        "   • Construction equipment tracking",
        "   • Cold chain / pharmaceutical logistics",
    ],
    [
        "● Primary Segments (Y1–Y2 Focus):",
        "   • Early adopters: tech-savvy urban Indians",
        "   • Pilot enterprises: 2–8 fleet units",
        "",
        "● Secondary Segments (Y3+):",
        "   • Mass retail via Amazon/online",
        "   • Mid-size fleets (30–250 units)",
        "",
        "● Key Insight:",
        "   B2C drives volume; B2B drives unit economics",
        "   and sticky SaaS revenue. Dual-segment strategy",
        "   de-risks the business.",
    ],
    "Who We Serve", "Segment Strategy"
)

# ═══════════════════ 2. VALUE PROPOSITION ═══════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_two_col_slide(slide, "2. Value Proposition",
    [
        "● For B2C Customers:",
        "   • Affordable: ₹3,000–₹4,500 vs AirTag ₹5,000+",
        "   • India-first: Local support, Indian maps",
        "   • AI-powered: Smart alerts, geofencing,",
        "     location history with AI Pro (₹999/yr)",
        "   • Family Plan (₹599/yr): Multi-device tracking",
        "     for the whole household",
        "   • Long battery life, water-resistant",
        "",
        "● For B2B Customers:",
        "   • Fleet tracking at 40–60% lower cost than",
        "     incumbents (₹4,999–₹5,999/unit)",
    ],
    [
        "   • SaaS console with real-time dashboards",
        "   • NB-IoT/LTE-M cellular connectivity",
        "   • Enterprise Hub + API for custom integration",
        "",
        "● Key Differentiators:",
        "   • Cost advantage: 25–50% below global brands",
        "   • Made for India: Built for Indian conditions,",
        "     roads, and temperatures",
        "   • Dual-revenue model: Hardware + Subscription",
        "     creates higher LTV than one-time sale",
        "   • Privacy-first: India-hosted data, compliant",
        "     with local regulations",
        "",
        "● Problem Solved:",
        "   \"Indians lose things daily — keys, wallets,",
        "   bags, vehicles. No affordable, India-ready",
        "   tracking solution exists.\"",
    ],
    "What We Offer", "Why TagX?"
)

# ═══════════════════ 3. CHANNELS ═══════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_two_col_slide(slide, "3. Channels",
    [
        "● Marketing Channels:",
        "   • Instagram & YouTube — demo videos,",
        "     influencer unboxings (B2C focus)",
        "   • Google Ads & Facebook Ads —",
        "     targeted acquisition (₹1L Y1 → ₹20L Y5)",
        "   • LinkedIn & industry events — B2B lead gen",
        "   • Referral program — users refer friends",
        "",
        "● Sales Channels:",
        "   • D2C website (Shopify/own store) — Y1–Y2",
        "   • Amazon India marketplace — scaling Y3+",
        "   • B2B direct sales team — from Y2 onwards",
    ],
    [
        "   • Enterprise pilots → eval → multi-year contract",
        "",
        "● Delivery Channels:",
        "   • Mobile app (iOS + Android) — primary",
        "     customer interface",
        "   • Web dashboard — for B2B fleet management",
        "   • Logistics partner — pan-India shipping",
        "   • Cloud platform (AWS/GCP) — real-time data",
        "",
        "● Channel Strategy:",
        "   Y1–Y2: Direct (website + founder-led sales)",
        "   Y3+: Marketplace + B2B team + retail expansion",
        "",
        "● Budget: ₹1L (Y1) → ₹20L (Y5)",
    ],
    "Reaching Customers", "Marketing → Sale → Delivery"
)

# ═══════════════════ 4. CUSTOMER RELATIONSHIPS ═══════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_two_col_slide(slide, "4. Customer Relationships",
    [
        "● B2C Relationships:",
        "   • Self-service app — setup, tracking,",
        "     subscription management",
        "   • In-app notifications & alerts — automated",
        "   • Email + chat support — ₹24K→₹8L/yr",
        "   • Social media community — Instagram,",
        "     WhatsApp groups for user tips",
        "   • Referral rewards — word-of-mouth growth",
        "",
        "● B2B Relationships:",
        "   • Dedicated account manager — from Y3",
        "   • Onboarding & training — enterprise pilots",
        "   • Quarterly business reviews — fleet analytics",
    ],
    [
        "   • API documentation & developer support",
        "",
        "● Retention Strategy:",
        "   • Subscription lock-in: Family & AI Pro plans",
        "     create recurring engagement",
        "   • Cloud sync: All data persists even if tag",
        "     is lost — user stays in ecosystem",
        "   • Regular firmware updates — OTA improvements",
        "",
        "● Support Scaling:",
        "   Y1: Founders handle all support",
        "   Y2: Part-time support hire",
        "   Y3+: Dedicated team (₹1.8L→₹8L/yr)",
        "",
        "● Goal: < 24hr response time, > 4.5★ app rating",
    ],
    "How We Engage", "Retention & Support"
)

# ═══════════════════ 5. REVENUE STREAMS ═══════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_two_col_slide(slide, "5. Revenue Streams",
    [
        "● Hardware Sales (60–75% of revenue):",
        "   • B2C Tags: ₹2,699–₹3,999/unit",
        "     250 → 10,000 units over 5 years",
        "   • B2B Fleet: ₹4,999–₹5,999/unit",
        "     5 → 200 enterprise units",
        "   • Enterprise Hub: ₹14,999/unit",
        "   • Accessories (straps, mounts): ~3–5% of tag rev",
        "",
        "● Subscription Revenue (15–25% of revenue):",
        "   • Family Plan: ₹599–₹999/yr, 15–45% attach",
        "   • AI Pro Plan: ₹999–₹1,500/yr, 15–40% attach",
        "   • AI inference cost: ₹100/user/yr → high margin",
        "   • B2B SaaS: ₹499–₹999/mo per subscriber",
    ],
    [
        "● Service Revenue (5–10% of revenue):",
        "   • Cellular data: ₹60–₹75/mo per device",
        "   • API subscriptions: ₹24,999/yr per enterprise",
        "   • White-label/OEM: ₹1L→₹15L by Y5",
        "",
        "● Revenue Breakdown (Y5 Est.):",
        "   • B2C Hardware:      ~₹4.0Cr  (70%)",
        "   • Subscriptions:      ~₹1.05Cr (18%)",
        "   • B2B:                ~₹0.43Cr (8%)",
        "   • Accessories:        ~₹0.20Cr (4%)",
        "   • Total:             ~₹5.68Cr",
        "",
        "● Key Metric:",
        "   ~3.0× YoY revenue growth",
        "   ₹8.3L (Y1) → ₹5.68Cr (Y5)",
    ],
    "How We Make Money", "Revenue Mix (Y5)"
)

# ═══════════════════ 6. KEY RESOURCES ═══════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_two_col_slide(slide, "6. Key Resources",
    [
        "● Physical:",
        "   • Office/co-working (Y1–Y2: WFH, Y3+ office)",
        "   • Test equipment & prototype lab",
        "   • Inventory warehouse (from Y3)",
        "",
        "● Intellectual:",
        "   • Hardware design (PCB, firmware, antenna)",
        "   • Mobile app (iOS + Android)",
        "   • Cloud platform & tracking algorithms",
        "   • Brand — TagX",
        "   • BIS certification (mandatory for India)",
        "   • Patents (filing in Y1, ₹1.5L)",
        "",
        "● Human:",
        "   • 2 founders (Y1) → 15 (Y5)",
        "   • Hardware engineers, firmware devs",
        "   • Full-stack developers (app + cloud)",
    ],
    [
        "   • B2B sales team (from Y3)",
        "   • Customer support team",
        "",
        "● Financial:",
        "   • Seed funding: ~₹42L (₹30L CapEx +",
        "     Y1–Y2 deficits + 10% contingency)",
        "   • Revenue reinvestment from Y3+",
        "",
        "● Partner Resources:",
        "   • China EMS manufacturing line",
        "   • Cloud infrastructure (AWS/GCP)",
        "   • IoT SIM/MVNO partnership",
        "",
        "● Resource Investment (Y1):",
        "   • Hardware R&D: ₹12L",
        "   • Tooling & molds: ₹10L",
        "   • App & platform: ₹6L",
        "   • BIS certification: ₹5L",
        "   • Office & equipment: ₹2.5L",
        "   • Legal & IP: ₹1.5L",
    ],
    "Assets We Need", "Investment Breakdown"
)

# ═══════════════════ 7. KEY ACTIVITIES ═══════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_two_col_slide(slide, "7. Key Activities",
    [
        "● Product R&D (Y1 Priority):",
        "   • PCB design & antenna tuning",
        "   • Firmware development (BLE + GPS + NB-IoT)",
        "   • Mobile app (React Native, iOS + Android)",
        "   • Cloud backend (real-time tracking API)",
        "   • Prototyping → pilot run → mass production",
        "   • Budget: ₹1L→₹3L/yr",
        "",
        "● Supply Chain & Manufacturing:",
        "   • China EMS partner selection & mgmt",
        "   • Component sourcing (GPS, BLE, battery)",
        "   • Quality control & BIS compliance",
        "   • Inventory planning & logistics",
        "",
        "● Marketing & Sales:",
        "   • B2C: Digital ads, content, influencer campaigns",
    ],
    [
        "   • B2B: Direct outreach, trade shows, pilots",
        "   • Budget: ₹1L→₹20L/yr",
        "",
        "● Platform Operations:",
        "   • Cloud uptime & scalability",
        "   • OTA firmware updates",
        "   • Security & data privacy compliance",
        "   • Budget: ₹30K→₹10L/yr",
        "",
        "● Customer Support:",
        "   • App-based ticketing system",
        "   • Email + phone support",
        "   • Self-service knowledge base",
        "   • Budget: ₹24K→₹8L/yr",
        "",
        "● Activity Focus by Year:",
        "   Y1: R&D + certification + pilot",
        "   Y2: Production scale + B2C marketing",
        "   Y3+: B2B expansion + platform scale",
    ],
    "What We Do Daily", "Activity Evolution"
)

# ═══════════════════ 8. KEY PARTNERSHIPS ═══════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_two_col_slide(slide, "8. Key Partnerships",
    [
        "● Manufacturing & Supply Chain:",
        "   • China EMS partner — PCB assembly,",
        "     final product assembly, quality control",
        "   • Component suppliers: GPS modules,",
        "     BLE chips, batteries, enclosures",
        "   • BIS certification labs — mandatory",
        "     India electronics certification",
        "",
        "● Technology Partners:",
        "   • Cloud: AWS or GCP — scalable hosting",
        "   • IoT connectivity: 1NCE / Hologram —",
        "     global NB-IoT/LTE-M MVNO at bulk rates",
        "   • AI inference: Groq API — low-cost",
        "     AI at ₹100/user/yr",
    ],
    [
        "● Go-to-Market Partners:",
        "   • Amazon India — primary e-commerce",
        "     channel from Y3 onwards",
        "   • Influencers & tech reviewers —",
        "     B2C awareness at low cost",
        "   • Fleet management software companies —",
        "     B2B referral partnerships",
        "",
        "● Strategic Value:",
        "   • Partnerships reduce capital needs:",
        "     EMS partner eliminates factory investment,",
        "     MVNO eliminates cellular infra build",
        "   • Partnerships add credibility:",
        "     AWS/Azure, BIS, Amazon — trust signals",
        "   • Partnerships accelerate growth:",
        "     Influencer reach, B2B referrals",
    ],
    "Who We Work With", "Why These Partners"
)

# ═══════════════════ 9. COST STRUCTURE ═══════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_two_col_slide(slide, "9. Cost Structure",
    [
        "● Variable Costs (COGS):",
        "   • BOM per B2C tag: ₹750→₹600 (China-sourced)",
        "   • Assembly: ₹80→₹60 per tag",
        "   • BOM per B2B unit: ₹1,000→₹800",
        "   • Cellular data: ₹15→₹12/dev/month",
        "   • Manufacturing overhead: 5–6% of COGS",
        "   • Total COGS (Y5): ~₹72L",
        "",
        "● Fixed Costs (OpEx):",
        "   • Team salaries: ₹1.8L→₹96L/yr",
        "     (founder stipend Y1 → market hires Y3+)",
        "   • Marketing: ₹1L→₹20L/yr",
        "   • Cloud infrastructure: ₹30K→₹10L/yr",
    ],
    [
        "   • AI inference: ₹100/premium user/yr",
        "   • R&D/prototyping: ₹1L→₹3L/yr",
        "   • B2B sales: ₹0→₹8L/yr",
        "   • Office: ₹0→₹4.8L/yr",
        "",
        "● Capital Expenditure:",
        "   • Initial CapEx: ₹30L (R&D, tooling,",
        "     certification, app, equipment, IP)",
        "   • Annual CapEx: ₹2L→₹8L/yr",
        "   • Depreciation: WDV 15% (hardware) +",
        "     SLM 3yr (software) + 5yr amortization",
        "",
        "● Cost Ratios (Y5):",
        "   • COGS as % of Rev: ~10% (scale benefit)",
        "   • OpEx as % of Rev: ~24%",
        "   • Gross Margin: ~90%",
        "   • EBITDA Margin: ~66%",
    ],
    "What It Costs", "Cost Efficiency"
)

# ═══════════════════ CLOSING SLIDE ═══════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)

add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
    "TagX — The Big Picture", font_size=44, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(2.5),
    "Y1    250 tags  |  ₹8.3L revenue  |  Founders only\n"
    "Y2    500 tags  |  ₹21.0L revenue  |  First hires\n"
    "Y3    1,300 tags  |  ₹65.3L revenue  |  Product-market fit\n"
    "Y4    3,000 tags  |  ₹1.78Cr revenue  |  BREAK-EVEN ✓\n"
    "Y5    10,000 tags  |  ₹5.68Cr revenue  |  Scaling\n\n"
    "Funding Required: ~₹40L  |  Cumulative Profit (Y5): ₹2.7Cr\n"
    "India GPS Tracker Market (2024): ₹1,214Cr  |  TagX Y5 Share: ~0.3%",
    font_size=20, color=RGBColor(0xC8, 0xE6, 0xC9), alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.6),
    "Built for India  |  Hardware + Subscription  |  B2C + B2B",
    font_size=18, color=RGBColor(0x81, 0xC7, 0x84), alignment=PP_ALIGN.CENTER)

# Save
prs.save("/Users/anoop/FY BTECH/Sem-4/TagX/TagX_Business_Model_Canvas.pptx")
print("✅ Saved to TagX_Business_Model_Canvas.pptx")
